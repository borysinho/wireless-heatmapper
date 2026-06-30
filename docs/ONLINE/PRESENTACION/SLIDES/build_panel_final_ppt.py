from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
ROOT = BASE.parents[3]
ASSETS = ROOT / "manual-usuario" / "assets"
IMG = BASE / "img"
OUT_DIR = BASE
EXPORT_DIR = BASE / "generated"
OUT = OUT_DIR / "WirelessHeatMapper-Presentacion-Final.pptx"


W = Inches(13.333)
H = Inches(7.5)


class C:
    NAVY = RGBColor(18, 32, 47)
    NAVY_2 = RGBColor(27, 42, 58)
    BLUE = RGBColor(41, 128, 185)
    BLUE_D = RGBColor(31, 97, 141)
    BLUE_L = RGBColor(235, 245, 251)
    CYAN = RGBColor(104, 196, 226)
    GREEN = RGBColor(24, 143, 103)
    AMBER = RGBColor(183, 121, 31)
    RED = RGBColor(194, 65, 59)
    INK = RGBColor(23, 32, 42)
    MUTED = RGBColor(93, 107, 122)
    LINE = RGBColor(217, 226, 236)
    BG = RGBColor(246, 248, 250)
    WHITE = RGBColor(255, 255, 255)
    YELLOW_L = RGBColor(255, 253, 231)


FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def emu(v: float) -> int:
    return int(v)


def svg_to_png(name: str, width: int = 1600) -> Path:
    src = ASSETS / f"{name}.svg"
    dst = EXPORT_DIR / f"{name}.png"
    EXPORT_DIR.mkdir(exist_ok=True)
    if not dst.exists() or dst.stat().st_mtime < src.stat().st_mtime:
        subprocess.run(
            [
                "inkscape",
                str(src),
                f"--export-filename={dst}",
                f"--export-width={width}",
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    return dst


def set_run(run, size=18, bold=False, color=C.INK, font=FONT_BODY):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    size=18,
    bold=False,
    color=C.INK,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_multiline(
    slide,
    lines: list[str],
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=C.INK,
    bullet=False,
    gap=6,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        if bullet:
            p.text = line
            p.font.name = FONT_BODY
            p.font.size = Pt(size)
            p.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = line
            set_run(run, size=size, color=color)
    return box


def fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def line(shape, color=C.LINE, width=1.2, transparency=0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_rect(slide, x, y, w, h, color=C.WHITE, border=C.LINE, radius=True, lw=1.2):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    fill(s, color)
    line(s, border, lw)
    return s


def add_pill(slide, text, x, y, w, h, color=C.BLUE, text_color=C.WHITE, size=12):
    s = add_rect(slide, x, y, w, h, color, color, radius=True, lw=0)
    tf = s.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=True, color=text_color)
    return s


def add_footer(slide, idx: int):
    add_textbox(
        slide,
        "Wireless HeatMapper · Grupo 24 · FICCT-UAGRM",
        Inches(0.55),
        Inches(7.08),
        Inches(6.2),
        Inches(0.22),
        size=8.5,
        color=C.MUTED,
    )
    add_textbox(
        slide,
        f"{idx:02d}/12",
        Inches(12.05),
        Inches(7.08),
        Inches(0.75),
        Inches(0.22),
        size=8.5,
        color=C.MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_title(slide, title: str, subtitle: str | None = None, idx: int | None = None):
    add_textbox(slide, title, Inches(0.58), Inches(0.35), Inches(8.7), Inches(0.46), size=25, bold=True, color=C.INK, font=FONT_HEAD)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.62), Inches(0.88), Inches(8.1), Inches(0.32), size=11.5, color=C.MUTED)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.24), Inches(1.05), Inches(0.04)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = C.BLUE
    slide.shapes[-1].line.fill.background()
    if idx:
        add_footer(slide, idx)


def add_card(slide, title: str, body: str, x, y, w, h, accent=C.BLUE, title_size=14, body_size=11.5):
    s = add_rect(slide, x, y, w, h, C.WHITE, C.LINE, True, 1.0)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = accent
    slide.shapes[-1].line.fill.background()
    add_textbox(slide, title, x + Inches(0.22), y + Inches(0.14), w - Inches(0.35), Inches(0.28), size=title_size, bold=True, color=C.INK)
    add_textbox(slide, body, x + Inches(0.22), y + Inches(0.52), w - Inches(0.35), h - Inches(0.62), size=body_size, color=C.MUTED)
    return s


def add_image_fit(slide, path: Path, x, y, w, h):
    pic = slide.shapes.add_picture(str(path), x, y)
    rw = w / pic.width
    rh = h / pic.height
    r = min(rw, rh)
    pic.width = int(pic.width * r)
    pic.height = int(pic.height * r)
    pic.left = int(x + (w - pic.width) / 2)
    pic.top = int(y + (h - pic.height) / 2)
    return pic


def connector(slide, x1, y1, x2, y2, color=C.BLUE):
    shp = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shp.line.color.rgb = color
    shp.line.width = Pt(2)
    return shp


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C.BG
    return slide


def slide_1(prs, logo):
    s = blank(prs)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H), C.NAVY)
    s.shapes[-1].line.fill.background()
    # Grilla sutil de fondo.
    for i in range(16):
        y = Inches(4.55 + i * 0.16)
        connector(s, Inches(0.0), y, Inches(13.333), y, RGBColor(45, 61, 77)).line.transparency = 60
    add_image_fit(s, logo, Inches(0.75), Inches(0.72), Inches(1.28), Inches(1.28))
    add_textbox(s, "Wireless HeatMapper", Inches(0.76), Inches(2.05), Inches(8.6), Inches(0.72), size=42, bold=True, color=C.WHITE, font=FONT_HEAD)
    add_textbox(s, "Sistema inteligente para análisis y optimización de cobertura WiFi", Inches(0.82), Inches(2.88), Inches(7.0), Inches(0.48), size=20, color=RGBColor(216, 232, 243))
    add_textbox(s, "Ingeniería de Software II · Grupo 24", Inches(0.82), Inches(4.85), Inches(5.0), Inches(0.35), size=15, bold=True, color=C.CYAN)
    add_textbox(s, "Cliente: Bulldog Tech.", Inches(0.82), Inches(5.28), Inches(3.4), Inches(0.28), size=13, color=RGBColor(216, 232, 243))
    add_textbox(s, "Jhasmany Jhunnior Fernandez Ortega · Herland Borys Quiroga Flores", Inches(0.82), Inches(5.67), Inches(7.0), Inches(0.28), size=12.5, color=RGBColor(216, 232, 243))
    # Línea cromática tipo heatmap.
    colors = [C.RED, C.AMBER, RGBColor(224, 205, 71), C.GREEN, C.BLUE]
    for i, col in enumerate(colors):
        shp = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.1 + i * 0.45), Inches(4.85), Inches(0.46), Inches(0.1))
        fill(shp, col)
        shp.line.fill.background()
    add_textbox(s, "FICCT-UAGRM · 2026", Inches(10.25), Inches(6.82), Inches(2.2), Inches(0.22), size=9.5, color=RGBColor(180, 198, 212), align=PP_ALIGN.RIGHT)


def slide_2(prs):
    s = blank(prs)
    add_title(s, "Problema y oportunidad", "El diseño WiFi interior no se valida bien solo con experiencia o estimaciones.", 2)
    add_card(s, "Entorno RF variable", "La señal cambia por distancia, obstáculos, interferencia y uso de canales.", Inches(0.75), Inches(1.75), Inches(3.85), Inches(1.25), C.BLUE)
    add_card(s, "Zonas débiles tardías", "Los problemas aparecen después de instalar o cuando el cliente ya reclama.", Inches(0.75), Inches(3.18), Inches(3.85), Inches(1.25), C.AMBER)
    add_card(s, "Decisión del cliente", "Se necesita evidencia visual y técnica para justificar una mejora.", Inches(0.75), Inches(4.61), Inches(3.85), Inches(1.25), C.GREEN)

    # Visual antes/después.
    add_rect(s, Inches(5.28), Inches(1.68), Inches(3.05), Inches(4.5), C.WHITE, C.LINE)
    add_textbox(s, "Antes", Inches(5.55), Inches(1.96), Inches(1.5), Inches(0.3), size=17, bold=True, color=C.INK)
    add_multiline(s, ["Estimación", "Experiencia técnica", "Resultados difíciles de explicar"], Inches(5.55), Inches(2.58), Inches(2.25), Inches(1.5), size=14, color=C.MUTED)
    for y in [3.95, 4.45, 4.95]:
        connector(s, Inches(5.65), Inches(y), Inches(7.75), Inches(y), RGBColor(190, 198, 207))

    add_rect(s, Inches(8.78), Inches(1.68), Inches(3.45), Inches(4.5), C.WHITE, C.LINE)
    add_textbox(s, "Con Wireless HeatMapper", Inches(9.05), Inches(1.96), Inches(2.85), Inches(0.3), size=17, bold=True, color=C.INK)
    for i, (label, col) in enumerate([("Medición RSSI", C.BLUE), ("Heatmap", C.GREEN), ("Propuesta IA", C.AMBER), ("Portal cliente", C.NAVY_2)]):
        add_pill(s, label, Inches(9.13), Inches(2.65 + i * 0.68), Inches(2.2), Inches(0.38), col, C.WHITE, 12)


def slide_3(prs, overview):
    s = blank(prs)
    add_title(s, "Solución propuesta", "Sistema integrado 100 % en línea para capturar, analizar, optimizar y publicar resultados WiFi.", 3)
    add_image_fit(s, overview, Inches(6.55), Inches(1.42), Inches(5.95), Inches(3.85))
    steps = [
        ("1", "Captura RSSI desde Android"),
        ("2", "Persistencia central en PostgreSQL"),
        ("3", "Heatmaps sobre planos calibrados"),
        ("4", "IA propone conjuntos AP derivados"),
        ("5", "Portal cliente por enlace"),
    ]
    for i, (num, text) in enumerate(steps):
        y = Inches(1.55 + i * 0.72)
        add_pill(s, num, Inches(0.78), y, Inches(0.42), Inches(0.42), C.BLUE, C.WHITE, 12)
        add_textbox(s, text, Inches(1.38), y + Inches(0.04), Inches(4.5), Inches(0.3), size=15.5, color=C.INK)
    add_textbox(s, "La app móvil es cliente delgado: captura y consulta; el estado de dominio vive en backend.", Inches(0.85), Inches(5.65), Inches(5.4), Inches(0.6), size=13.5, color=C.MUTED)


def slide_4(prs):
    s = blank(prs)
    add_title(s, "Arquitectura online", "Una sola fuente de verdad: el backend central y PostgreSQL.", 4)
    blocks = [
        ("App móvil\nFlutter / Dart", Inches(0.8), Inches(2.15), C.BLUE_L, C.BLUE),
        ("Web admin + portal\nReact + TypeScript", Inches(0.8), Inches(4.15), C.BLUE_L, C.BLUE),
        ("Nginx\nreverse proxy", Inches(4.1), Inches(3.15), C.WHITE, C.NAVY_2),
        ("Backend REST + IA\nFastAPI / Python", Inches(6.45), Inches(3.15), C.WHITE, C.BLUE),
        ("PostgreSQL 15\nfuente de verdad", Inches(9.45), Inches(3.15), RGBColor(234, 247, 234), C.GREEN),
    ]
    for text, x, y, bg, accent in blocks:
        add_rect(s, x, y, Inches(2.4), Inches(0.96), bg, accent, True, 1.6)
        add_textbox(s, text, x + Inches(0.15), y + Inches(0.14), Inches(2.1), Inches(0.58), size=13.2, bold=True, color=C.INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    connector(s, Inches(3.2), Inches(2.63), Inches(4.1), Inches(3.35))
    connector(s, Inches(3.2), Inches(4.63), Inches(4.1), Inches(3.85))
    connector(s, Inches(6.5), Inches(3.63), Inches(6.45), Inches(3.63))
    connector(s, Inches(8.85), Inches(3.63), Inches(9.45), Inches(3.63))
    add_card(s, "Decisiones clave", "Sin sqflite/drift · Sin sincronización diferida · Operación REST en línea · Cálculo pesado en backend", Inches(0.85), Inches(5.78), Inches(11.2), Inches(0.72), C.BLUE, 13.5, 11.5)


def slide_5(prs):
    s = blank(prs)
    add_title(s, "Alcance y Scrum", "Desarrollo incremental con trazabilidad entre RP, HU y sprints.", 5)
    items = [
        ("S0-S1", "Base técnica, admin,\nauth y proyectos", C.BLUE),
        ("S2-S3", "Planos, calibración,\ncaptura y puntos", C.GREEN),
        ("S4", "Conjuntos AP y\nheatmaps backend", C.AMBER),
        ("S5", "IA RF y comparación\nde propuestas", C.RED),
        ("S6", "Portal cliente y\nenlaces únicos", C.NAVY_2),
    ]
    x0 = Inches(0.78)
    for i, (sprint, value, col) in enumerate(items):
        x = x0 + Inches(i * 2.45)
        add_pill(s, sprint, x, Inches(2.0), Inches(1.25), Inches(0.42), col, C.WHITE, 13)
        add_rect(s, x, Inches(2.62), Inches(2.0), Inches(1.23), C.WHITE, C.LINE)
        add_textbox(s, value, x + Inches(0.13), Inches(2.88), Inches(1.72), Inches(0.65), size=12.2, bold=True, color=C.INK, align=PP_ALIGN.CENTER)
        if i < len(items) - 1:
            connector(s, x + Inches(2.03), Inches(3.24), x + Inches(2.35), Inches(3.24), C.BLUE)
    add_card(s, "Dato clave", "17 historias de usuario vigentes · 142 puntos de historia · PB-14 eliminada por modalidad online", Inches(1.05), Inches(5.0), Inches(11.2), Inches(0.95), C.BLUE, 16, 14)


def slide_6(prs, mobile):
    s = blank(prs)
    add_title(s, "Flujo de uso completo", "Del trabajo de campo a la entrega controlada al cliente.", 6)
    add_image_fit(s, mobile, Inches(7.15), Inches(1.25), Inches(5.55), Inches(4.85))
    steps = [
        "Admin registra técnicos y clientes",
        "Técnico crea proyecto y sube plano",
        "Calibra escala y marca puntos",
        "App envía mediciones RSSI",
        "Backend genera heatmap",
        "Web genera propuesta IA",
        "Admin publica contenido",
        "Cliente accede por enlace",
    ]
    for i, text in enumerate(steps):
        y = Inches(1.45 + i * 0.48)
        add_pill(s, str(i + 1), Inches(0.78), y, Inches(0.34), Inches(0.34), C.BLUE if i < 5 else C.GREEN, C.WHITE, 10.5)
        add_textbox(s, text, Inches(1.28), y + Inches(0.03), Inches(5.2), Inches(0.25), size=13.2, color=C.INK)


def entity(slide, title, fields, x, y, w=1.95, h=0.88, color=C.BLUE):
    add_rect(slide, x, y, Inches(w), Inches(h), C.WHITE, color, True, 1.4)
    add_textbox(slide, title, x + Inches(0.1), y + Inches(0.1), Inches(w - 0.2), Inches(0.22), size=11.5, bold=True, color=C.INK, align=PP_ALIGN.CENTER)
    add_textbox(slide, fields, x + Inches(0.12), y + Inches(0.38), Inches(w - 0.24), Inches(0.35), size=8.6, color=C.MUTED, align=PP_ALIGN.CENTER)


def slide_7(prs):
    s = blank(prs)
    add_title(s, "Núcleo del modelo de datos RF", "Mostrar solo la parte relacional que sostiene medición, conjuntos AP y heatmaps.", 7)
    # Entidades principales del núcleo RF.
    entity(s, "proyecto", "cliente · estado\ntecnico_id", Inches(0.75), Inches(1.78), 1.8, 0.84, C.BLUE)
    entity(s, "plano", "escala · poligono\nancho/alto", Inches(3.1), Inches(1.78), 1.8, 0.84, C.BLUE)
    entity(s, "punto_medicion", "pos_x · pos_y\nnivel", Inches(5.45), Inches(1.78), 1.8, 0.84, C.GREEN)
    entity(s, "lectura_rssi", "SSID · BSSID\ncanal · RSSI", Inches(7.85), Inches(1.78), 1.8, 0.84, C.GREEN)
    entity(s, "conjunto_ap", "origen · proposito\nconjunto_origen_id", Inches(3.1), Inches(4.0), 1.9, 0.88, C.AMBER)
    entity(s, "conjunto_ap_item", "BSSID · posicion\naccion IA", Inches(5.58), Inches(4.0), 1.95, 0.88, C.AMBER)
    entity(s, "mapa_calor", "matriz · escala\nfirma_mediciones", Inches(8.15), Inches(4.0), 1.95, 0.88, C.RED)
    # Relaciones principales.
    connector(s, Inches(2.55), Inches(2.2), Inches(3.1), Inches(2.2))
    connector(s, Inches(4.9), Inches(2.2), Inches(5.45), Inches(2.2))
    connector(s, Inches(7.25), Inches(2.2), Inches(7.85), Inches(2.2))
    connector(s, Inches(4.0), Inches(2.62), Inches(4.0), Inches(4.0), C.AMBER)
    connector(s, Inches(5.0), Inches(4.44), Inches(5.58), Inches(4.44), C.AMBER)
    connector(s, Inches(7.53), Inches(4.44), Inches(8.15), Inches(4.44), C.RED)
    connector(s, Inches(4.88), Inches(2.2), Inches(8.15), Inches(4.05), C.RED).line.transparency = 25
    add_card(
        s,
        "Lectura técnica",
        "La IA se persiste como conjunto_ap de origen 'ia'. Se omiten login, tokens y entidades administrativas porque no pertenecen al núcleo RF.",
        Inches(0.95),
        Inches(5.62),
        Inches(10.9),
        Inches(0.88),
        C.BLUE,
        12.8,
        10.4,
    )


def slide_8(prs, heatmap):
    s = blank(prs)
    add_title(s, "Heatmaps y criterios RF", "Mapas generados desde evidencia RSSI real y criterios técnicos explícitos.", 8)
    add_image_fit(s, heatmap, Inches(6.9), Inches(1.25), Inches(5.45), Inches(4.52))
    facts = [
        ("Fuente observada", "lectura_rssi.origen = CAMPO", C.BLUE),
        ("Algoritmo persistido", "IDW", C.GREEN),
        ("Objetivo de diseño", "RSSI >= -70 dBm", C.AMBER),
        ("Zona muerta", "RSSI < -90 dBm", C.RED),
        ("Android >= 8.0", "4 scans / 2 min", C.NAVY_2),
    ]
    for i, (t, b, col) in enumerate(facts):
        y = Inches(1.35 + i * 0.9)
        add_rect(s, Inches(0.78), y, Inches(5.35), Inches(0.7), C.WHITE, C.LINE)
        slide_bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.78), y, Inches(0.08), Inches(0.7))
        fill(slide_bar, col)
        slide_bar.line.fill.background()
        add_textbox(s, t, Inches(1.05), y + Inches(0.12), Inches(2.3), Inches(0.22), size=12.6, bold=True, color=C.INK)
        add_textbox(s, b, Inches(3.35), y + Inches(0.13), Inches(2.45), Inches(0.22), size=12.3, color=C.MUTED)
    add_textbox(s, "IDW interpola mediciones reales; no simula una red inexistente.", Inches(0.86), Inches(6.12), Inches(5.7), Inches(0.32), size=12.5, bold=True, color=C.BLUE_D)


def slide_9(prs):
    s = blank(prs)
    add_title(s, "IA RF y optimización AP", "La IA propone escenarios futuros sin alterar la evidencia real capturada.", 9)
    nodes = [
        ("Conjunto técnico", "APs relevantes\nmedidos en campo", C.BLUE),
        ("Predictor RF", "FSPL/log-distance\ncalibración local", C.AMBER),
        ("Conjunto IA", "APs propuestos\norigen = ia", C.GREEN),
        ("Heatmap proyectado", "IDW sobre lecturas\nIA_ESTIMADA", C.RED),
    ]
    x0 = Inches(0.8)
    for i, (title, body, col) in enumerate(nodes):
        x = x0 + Inches(i * 3.03)
        add_rect(s, x, Inches(2.2), Inches(2.25), Inches(1.35), C.WHITE, col, True, 1.6)
        add_textbox(s, title, x + Inches(0.16), Inches(2.42), Inches(1.92), Inches(0.24), size=13.2, bold=True, color=C.INK, align=PP_ALIGN.CENTER)
        add_textbox(s, body, x + Inches(0.16), Inches(2.82), Inches(1.92), Inches(0.42), size=10.8, color=C.MUTED, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            connector(s, x + Inches(2.25), Inches(2.88), x + Inches(2.85), Inches(2.88), C.BLUE)
    add_card(s, "Regla de defensa", "IDW y FSPL no compiten: IDW persiste mapas; FSPL estima escenarios futuros para alimentar la IA.", Inches(1.15), Inches(4.72), Inches(10.8), Inches(0.85), C.BLUE, 15.5, 13.5)


def slide_10(prs, portal):
    s = blank(prs)
    add_title(s, "Seguridad y portal cliente", "El cliente ve resultados interactivos sin acceso interno al sistema.", 10)
    add_image_fit(s, portal, Inches(7.0), Inches(1.28), Inches(5.2), Inches(4.6))
    add_card(s, "Usuarios internos", "JWT + refresh token + bcrypt\nRoles para administración y operación técnica", Inches(0.8), Inches(1.55), Inches(5.35), Inches(1.05), C.BLUE, 14, 12)
    add_card(s, "Portal cliente", "/portal/:token\nEnlace con expiración y revocación", Inches(0.8), Inches(2.9), Inches(5.35), Inches(1.05), C.GREEN, 14, 12)
    add_card(s, "Visibilidad restringida", "Contenido limitado a conjunto_ids y mapa_ids\nSin exposición de otros proyectos", Inches(0.8), Inches(4.25), Inches(5.35), Inches(1.05), C.AMBER, 14, 12)
    add_textbox(s, "El portal reemplaza el PDF por una entrega visual, interactiva y controlada.", Inches(0.9), Inches(6.05), Inches(5.4), Inches(0.35), size=12.5, bold=True, color=C.BLUE_D)


def slide_11(prs):
    s = blank(prs)
    add_title(s, "Calidad y evidencia de implementación", "Solución implementada por capas y acompañada por pruebas automatizadas.", 11)
    add_card(s, "Backend", "Routers: auth, usuarios, clientes, proyectos, planos, mediciones, heatmaps, escenarios y share.\nTests pytest para endpoints, heatmaps, IA y portal.", Inches(0.75), Inches(1.72), Inches(3.75), Inches(2.25), C.BLUE, 15, 11.5)
    add_card(s, "Web", "React + TypeScript para admin y portal.\nVitest / React Testing Library en portal y utilidades API.", Inches(4.82), Inches(1.72), Inches(3.75), Inches(2.25), C.GREEN, 15, 11.5)
    add_card(s, "Móvil", "Flutter con datasources remotos.\nTests unit/widget para auth, proyectos, planos, captura y heatmap.", Inches(8.9), Inches(1.72), Inches(3.75), Inches(2.25), C.AMBER, 15, 11.5)
    add_rect(s, Inches(1.35), Inches(5.05), Inches(10.7), Inches(0.82), C.WHITE, C.LINE)
    add_textbox(s, "Infraestructura: migraciones Alembic · Docker Compose · Nginx · PostgreSQL 15", Inches(1.55), Inches(5.28), Inches(10.25), Inches(0.28), size=15, bold=True, color=C.INK, align=PP_ALIGN.CENTER)


def slide_12(prs):
    s = blank(prs)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H), C.NAVY)
    s.shapes[-1].line.fill.background()
    add_textbox(s, "Valor entregado", Inches(0.78), Inches(0.62), Inches(5.2), Inches(0.6), size=34, bold=True, color=C.WHITE, font=FONT_HEAD)
    add_textbox(s, "Wireless HeatMapper convierte el relevamiento WiFi en un proceso medible, trazable y compartible.", Inches(0.82), Inches(1.35), Inches(8.9), Inches(0.52), size=18, color=RGBColor(216, 232, 243))
    cards = [
        ("Evidencia centralizada", "PostgreSQL como fuente única", C.BLUE),
        ("Heatmaps reales", "RSSI observado + IDW", C.GREEN),
        ("IA trazable", "Conjuntos derivados desde campo", C.AMBER),
        ("Entrega profesional", "Portal seguro por enlace", C.CYAN),
    ]
    for i, (title, body, col) in enumerate(cards):
        x = Inches(0.9 + i * 3.05)
        add_rect(s, x, Inches(3.1), Inches(2.55), Inches(1.34), RGBColor(31, 48, 65), RGBColor(72, 91, 109), True, 1.1)
        add_pill(s, "", x + Inches(0.2), Inches(3.35), Inches(0.46), Inches(0.12), col, col, 1)
        add_textbox(s, title, x + Inches(0.2), Inches(3.62), Inches(2.1), Inches(0.28), size=14.5, bold=True, color=C.WHITE)
        add_textbox(s, body, x + Inches(0.2), Inches(4.05), Inches(2.05), Inches(0.38), size=11.5, color=RGBColor(195, 211, 224))
    add_textbox(s, "Cadena completa de evidencia técnica para tomar decisiones WiFi.", Inches(0.95), Inches(6.25), Inches(11.1), Inches(0.38), size=18, bold=True, color=C.CYAN, align=PP_ALIGN.CENTER)
    add_textbox(s, "12/12", Inches(12.05), Inches(7.08), Inches(0.75), Inches(0.22), size=8.5, color=RGBColor(180, 198, 212), align=PP_ALIGN.RIGHT)


def main():
    overview = svg_to_png("vista-general", 1700)
    mobile = svg_to_png("app-movil", 1700)
    heatmap = svg_to_png("heatmap-ia", 1700)
    portal = svg_to_png("portal-cliente", 1700)
    logo = IMG / "logo.png"

    prs = new_deck()
    slide_1(prs, logo)
    slide_2(prs)
    slide_3(prs, overview)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs, mobile)
    slide_7(prs)
    slide_8(prs, heatmap)
    slide_9(prs)
    slide_10(prs, portal)
    slide_11(prs)
    slide_12(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
